"""
generate_ppt.py
基于 report_template.pptx 模板，读取 CSV 数据并生成最终 PPT 报告。

模板包含 3 张预设幻灯片（索引 0=封面, 1=分隔页, 2=数据页），
本脚本克隆模板幻灯片，填入文字和表格数据后删除原始模板页。
"""
import csv
import os
import copy
import zipfile
import lxml.etree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── 配置 ──────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(SCRIPT_DIR, "report_template.pptx")
COMPANY_DIRS = ["company1", "company2", "company3"]

PRISM_BLUE = RGBColor(0x00, 0x70, 0xC0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x00, 0x70, 0xC0)
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF0, 0xFE)

csv_files = [
    ("employees.csv", "Employee Information & Payroll",
     "Employee personal details, salary structure and deductions"),
    ("products.csv", "Product Inventory Management",
     "Product catalog, pricing, stock levels and warehouse locations"),
    ("orders.csv", "Order Records",
     "Customer orders with payment, shipping and fulfillment status"),
]


# ── helpers ────────────────────────────────────────────
def _get_price_comparison_slide_info(report_path):
    """从已生成的 report.pptx 中获取 \"Monthly Price Comparison\" 数据页的
    (1-based slide_number, sldId, slide_title)。

    返回 (slide_number, sld_id, title) 或 (None, None, None)。
    """
    SECTION_TITLE_PREFIX = "Monthly Price Comparison"

    prs = Presentation(report_path)

    # 从 presentation.xml 中读取各幻灯片的内部 sldId
    with zipfile.ZipFile(report_path, "r") as z:
        pres_xml = z.read("ppt/presentation.xml")
    pres_root = ET.fromstring(pres_xml)
    sld_id_lst = pres_root.find(".//" + qn("p:sldIdLst"))
    sld_ids = [int(e.get("id")) for e in sld_id_lst] if sld_id_lst is not None else []

    # 查找第一个包含 PageInfo shape 的 "Monthly Price Comparison" 幻灯片（即数据页而非分隔页）
    for idx, slide in enumerate(prs.slides):
        has_section_title = False
        has_page_info = False
        title_text = ""
        for sp in slide.shapes:
            if sp.name == "SectionTitle" and hasattr(sp, "text"):
                text = sp.text.strip()
                if text.startswith(SECTION_TITLE_PREFIX):
                    has_section_title = True
                    title_text = text
            if sp.name == "PageInfo":
                has_page_info = True
        if has_section_title and has_page_info:
            slide_number = idx + 1
            sld_id = sld_ids[idx] if idx < len(sld_ids) else 256 + idx
            return slide_number, sld_id, title_text

    return None, None, None


def _find_shape(slide, name):
    """Return shape by name, or None."""
    for sp in slide.shapes:
        if sp.name == name:
            return sp
    return None


def _set_text(shape, text, size=None, bold=None, color=None, align=None):
    """Replace text while preserving template formatting when no overrides given."""
    tf = shape.text_frame
    tf.word_wrap = True

    if tf.paragraphs and tf.paragraphs[0].runs:
        # shape already has a run with template formatting — keep it
        run = tf.paragraphs[0].runs[0]
        run.text = text
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        if align is not None:
            tf.paragraphs[0].alignment = align
    else:
        # no existing run, create fresh
        tf.clear()
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = text
        if size:
            r.font.size = Pt(size)
        if bold is not None:
            r.font.bold = bold
        if color:
            r.font.color.rgb = color
        r.font.name = "Calibri"


def _read_cell_style(cell):
    """从模板表格单元格中提取字体样式。"""
    style = {}
    try:
        tf = cell.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            font = tf.paragraphs[0].runs[0].font
            if font.size is not None:
                style['font_size'] = font.size
            if font.name is not None:
                style['font_name'] = font.name
            if font.color and font.color.type is not None:
                style['font_color'] = font.color.rgb
            if font.bold is not None:
                style['bold'] = font.bold
    except Exception:
        pass
    try:
        p = cell.text_frame.paragraphs[0]
        if p.alignment is not None:
            style['align'] = p.alignment
    except Exception:
        pass
    return style


def _get_cell_fill_rgb(cell):
    """从单元格 XML 读取纯色填充色，返回 RGBColor 或 None。"""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        return None
    solidFill = tcPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is not None and srgbClr.get('val'):
        return RGBColor.from_string(srgbClr.get('val'))
    return None


def _set_cell(cell, text, style=None):
    """写入单元格文本并应用模板样式，缺少的属性用硬编码默认值补齐。"""
    style = style or {}
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    if 'align' in style and style['align'] is not None:
        p.alignment = style['align']
    r = p.add_run()
    r.text = str(text)
    if style.get('font_size'):
        r.font.size = style['font_size']
    if style.get('bold') is not None:
        r.font.bold = style['bold']
    if style.get('font_color'):
        r.font.color.rgb = style['font_color']
    r.font.name = style.get('font_name') or "Calibri"
    cell.vertical_anchor = 1  # middle


def _fill_cell(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _clone_slide(prs, template_slide, exclude_names=None):
    """Deep-copy a template slide into the presentation and return the new slide.
    exclude_names: list of shape names to skip (e.g. 'SampleTable').
    """
    exclude = set(exclude_names or [])
    slide_layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(slide_layout)

    # remove default elements from new slide's spTree
    sp_tree = new_slide.shapes._spTree
    for elem in list(sp_tree):
        if elem.tag.endswith("}sp") or elem.tag.endswith("}graphicFrame"):
            sp_tree.remove(elem)

    # copy all shape elements from template, skipping excluded ones
    src_tree = template_slide.shapes._spTree
    for elem in src_tree:
        if elem.tag.endswith("}sp") or elem.tag.endswith("}graphicFrame") or elem.tag.endswith("}pic"):
            # check if this element's name is in the exclude list
            nvSpPr = elem.find(qn("p:nvSpPr"))
            if nvSpPr is not None:
                nvPr = nvSpPr.find(qn("p:cNvPr"))
                if nvPr is not None and nvPr.get("name") in exclude:
                    continue
            # for graphicFrame (tables), check nvGrpSpPr/cNvPr
            nvGrpFr = elem.find(qn("p:nvGraphicFramePr"))
            if nvGrpFr is not None:
                nvPr = nvGrpFr.find(qn("p:cNvPr"))
                if nvPr is not None and nvPr.get("name") in exclude:
                    continue
            sp_tree.append(copy.deepcopy(elem))

    return new_slide


def read_csv(filepath):
    with open(filepath, newline="", encoding="utf-8") as f:
        reader = csv.reader(f)
        headers = next(reader)
        data = [row for row in reader]
    return headers, data


# ── slide builders ────────────────────────────────────
def build_cover(prs, template_slide):
    slide = _clone_slide(prs, template_slide)

    sh = _find_shape(slide, "Title")
    if sh:
        _set_text(sh, "CSV Data Report")

    sh = _find_shape(slide, "Subtitle")
    if sh:
        _set_text(sh, "Overview of three datasets from CSV files")

    sh = _find_shape(slide, "ContentList")
    if sh:
        tf = sh.text_frame
        tf.word_wrap = True
        # preserve the template run's font as base style
        base_font = None
        if tf.paragraphs and tf.paragraphs[0].runs:
            base_font = tf.paragraphs[0].runs[0].font
        tf.clear()

        for i, (_, title, desc) in enumerate(csv_files):
            if i > 0:
                tf.add_paragraph().space_before = Pt(6)

            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            rn = p.add_run()
            rn.text = f"{i + 1}. "
            rn.font.size = Pt(18)
            rn.font.bold = True
            if base_font and base_font.color and base_font.color.rgb:
                rn.font.color.rgb = base_font.color.rgb
            else:
                rn.font.color.rgb = PRISM_BLUE
            rn.font.name = base_font.name if base_font and base_font.name else "Calibri"

            rt = p.add_run()
            rt.text = title
            rt.font.size = Pt(18)
            rt.font.bold = True
            if base_font and base_font.color and base_font.color.rgb:
                # slightly darker than accent for title text
                rt.font.color.rgb = DARK_GRAY
            else:
                rt.font.color.rgb = DARK_GRAY
            rt.font.name = base_font.name if base_font and base_font.name else "Calibri"

            pd = tf.add_paragraph()
            pd.alignment = PP_ALIGN.LEFT
            rd = pd.add_run()
            rd.text = f"     {desc}"
            rd.font.size = Pt(14)
            if base_font and base_font.color and base_font.color.rgb:
                rd.font.color.rgb = MEDIUM_GRAY
            else:
                rd.font.color.rgb = MEDIUM_GRAY
            rd.font.name = base_font.name if base_font and base_font.name else "Calibri"


def build_section(prs, template_slide, title, rows, cols):
    slide = _clone_slide(prs, template_slide)

    sh = _find_shape(slide, "SectionTitle")
    if sh:
        _set_text(sh, title, align=PP_ALIGN.CENTER)

    sh = _find_shape(slide, "DimensionInfo")
    if sh:
        _set_text(sh, f"{rows} rows  x  {cols} columns", align=PP_ALIGN.CENTER)


def build_table(prs, template_slide, headers, data_chunk,
                page_num, total_pages, section_title, row_height, max_table_height,
                header_style=None, data_style=None, header_fill=None, alt_fill=None):
    num_cols = len(headers)
    num_rows = len(data_chunk) + 1

    # clone slide but skip the sample table
    slide = _clone_slide(prs, template_slide, exclude_names=["SampleTable"])

    sh = _find_shape(slide, "SectionTitle")
    if sh:
        _set_text(sh, section_title)

    sh = _find_shape(slide, "PageInfo")
    if sh:
        _set_text(sh, f"Page {page_num} / {total_pages}", align=PP_ALIGN.RIGHT)

    # calculate actual height: row_height * num_rows, capped at max
    TABLE_TOP = Inches(0.85)
    table_height = min(row_height * num_rows, max_table_height)

    # add data table
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), TABLE_TOP,
        Inches(12.3), table_height,
    )
    tbl = tbl_shape.table
    col_w = int(Inches(12.3) / num_cols)
    for i in range(num_cols):
        tbl.columns[i].width = col_w

    # header
    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        _set_cell(c, h, style=header_style)
        if header_fill:
            _fill_cell(c, header_fill)

    # data
    for ri, row in enumerate(data_chunk):
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            _set_cell(c, val, style=data_style)
            if ri % 2 == 1 and alt_fill:
                _fill_cell(c, alt_fill)

    return tbl_shape


# ── main ───────────────────────────────────────────────
def _read_table_config(template_slide):
    """Read table config + cell styles from the sample table in DATA_TABLE template slide.
    Returns (data_rows_per_page, single_row_height, max_table_height,
             header_style, data_style, header_fill, alt_fill).
    """
    for sp in template_slide.shapes:
        if sp.name == "SampleTable" and sp.has_table:
            total_rows = len(sp.table.rows)
            data_rows = total_rows - 1
            table_height = sp.height
            row_height = table_height // total_rows

            header_style = _read_cell_style(sp.table.cell(0, 0))
            header_fill = _get_cell_fill_rgb(sp.table.cell(0, 0))
            data_style = _read_cell_style(sp.table.cell(1, 0))
            alt_fill = None
            if total_rows > 2:
                alt_fill = _get_cell_fill_rgb(sp.table.cell(2, 0))

            # 模板未显式设置的属性用代码默认值补齐
            _DH = {'font_size': Pt(10), 'font_name': 'Calibri',
                   'font_color': TABLE_HEADER_FG, 'bold': True, 'align': PP_ALIGN.CENTER}
            _DD = {'font_size': Pt(9), 'font_name': 'Calibri',
                   'font_color': DARK_GRAY, 'bold': False, 'align': PP_ALIGN.LEFT}
            for k, v in _DH.items():
                header_style.setdefault(k, v)
            for k, v in _DD.items():
                data_style.setdefault(k, v)
            header_fill = header_fill or TABLE_HEADER_BG
            alt_fill = alt_fill or TABLE_ROW_ALT

            print(f"  SampleTable: {total_rows} rows, height={table_height} EMU")
            print(f"  -> {data_rows} data rows/page, {row_height} EMU/row")
            print(f"  header style: {header_style}, fill: {header_fill}")
            print(f"  data style: {data_style}, alt fill: {alt_fill}")
            return data_rows, row_height, table_height, header_style, data_style, header_fill, alt_fill

    # fallback
    print("  Warning: SampleTable not found, defaulting to 5 rows/page")
    default_h = Inches(6.35)
    hs = {'font_size': Pt(10), 'font_name': 'Calibri', 'font_color': TABLE_HEADER_FG,
          'bold': True, 'align': PP_ALIGN.CENTER}
    ds = {'font_size': Pt(9), 'font_name': 'Calibri', 'font_color': DARK_GRAY,
          'bold': False, 'align': PP_ALIGN.LEFT}
    return 5, default_h // 6, default_h, hs, ds, TABLE_HEADER_BG, TABLE_ROW_ALT


def compare_order_prices(filepath1, filepath2):
    """对比两个订单CSV的单价，返回 |差价| > 5 的产品列表及带符号差价。
    Returns (headers, data, diffs)。
    """
    headers1, data1 = read_csv(filepath1)
    headers2, data2 = read_csv(filepath2)

    try:
        name_idx1 = headers1.index("ProductName")
        price_idx1 = headers1.index("UnitPrice")
        name_idx2 = headers2.index("ProductName")
        price_idx2 = headers2.index("UnitPrice")
    except ValueError:
        return None, None, None

    price_map1 = {row[name_idx1]: float(row[price_idx1]) for row in data1}
    price_map2 = {row[name_idx2]: float(row[price_idx2]) for row in data2}

    common = sorted(set(price_map1.keys()) & set(price_map2.keys()))

    result_headers = ["ProductName", "Month1 Price (¥)", "Month2 Price (¥)", "Difference (¥)"]
    result_data = []
    diffs = []
    for prod in common:
        p1 = price_map1[prod]
        p2 = price_map2[prod]
        diff = p2 - p1
        if abs(diff) > 5:
            prefix = "+" if diff > 0 else ""
            result_data.append([prod, f"{p1:.0f}", f"{p2:.0f}", f"{prefix}{diff:.0f}"])
            diffs.append(diff)

    return result_headers, result_data, diffs


def generate_company_report(company_dir):
    """为单个公司目录生成 report.pptx。"""
    data_dir = os.path.join(SCRIPT_DIR, company_dir)
    output_path = os.path.join(data_dir, "report.pptx")

    prs = Presentation(TEMPLATE)

    # grab the 3 template slides (will be deleted later)
    tmpl_cover   = prs.slides[0]
    tmpl_section = prs.slides[1]
    tmpl_table   = prs.slides[2]

    # read table config from template's sample table
    rows_per_page, row_height, max_table_height, \
        header_style, data_style, header_fill, alt_fill = _read_table_config(tmpl_table)

    # ── build cover ───────────────────────────────────
    build_cover(prs, tmpl_cover)

    # ── build each CSV section ────────────────────────
    for filename, title, _ in csv_files:
        filepath = os.path.join(data_dir, filename)
        headers, data = read_csv(filepath)
        total_rows = len(data)
        total_pages = (total_rows + rows_per_page - 1) // rows_per_page

        build_section(prs, tmpl_section, title, total_rows, len(headers))

        for page in range(total_pages):
            start = page * rows_per_page
            end = min(start + rows_per_page, total_rows)
            build_table(prs, tmpl_table, headers, data[start:end],
                        page + 1, total_pages, title, row_height, max_table_height,
                        header_style, data_style, header_fill, alt_fill)

    # ── price comparison: orders vs orders2 ───────────
    order1_path = os.path.join(data_dir, "orders.csv")
    order2_path = os.path.join(data_dir, "orders2.csv")
    comp_headers, comp_data, comp_diffs = compare_order_prices(order1_path, order2_path)

    if comp_headers and comp_data:
        section_title = "Monthly Price Comparison (|diff| > 5)"
        build_section(prs, tmpl_section, section_title, len(comp_data), len(comp_headers))

        total_pages = (len(comp_data) + rows_per_page - 1) // rows_per_page
        for page in range(total_pages):
            start = page * rows_per_page
            end = min(start + rows_per_page, len(comp_data))
            tbl_shape = build_table(prs, tmpl_table, comp_headers, comp_data[start:end],
                                     page + 1, total_pages, section_title, row_height, max_table_height,
                                     header_style, data_style, header_fill, alt_fill)
            # 对 Difference 列（第4列，索引3）按差价区间着色
            tbl = tbl_shape.table
            YELLOW = RGBColor(0xFF, 0xFF, 0x00)
            RED = RGBColor(0xFF, 0x00, 0x00)
            GREEN = RGBColor(0x00, 0x80, 0x00)
            for ri, diff_val in enumerate(comp_diffs[start:end]):
                cell = tbl.cell(ri + 1, 3)
                if diff_val > 0:
                    if diff_val >= 100:
                        _fill_cell(cell, RED)
                    else:
                        _fill_cell(cell, YELLOW)
                else:
                    _fill_cell(cell, GREEN)
    else:
        print("  No price differences >5 found between orders.csv and orders2.csv")

    # ── remove the 3 original template slides ─────────
    # delete in reverse order to keep indices stable
    for idx in [2, 1, 0]:
        rId = prs.slides._sldIdLst[idx].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

    prs.save(output_path)
    print(f"  Done -> {output_path}")
    print(f"  Total slides: {len(prs.slides)}")


def generate_summary_report():
    """在根目录生成总结性 PPT，统计各公司价格差异的绝对值数量分布（每100为区间）。"""
    output_path = os.path.join(SCRIPT_DIR, "summary_report.pptx")

    prs = Presentation(TEMPLATE)
    tmpl_cover = prs.slides[0]
    tmpl_section = prs.slides[1]
    tmpl_table = prs.slides[2]

    _, row_height, max_table_height, \
        header_style, data_style, header_fill, alt_fill = _read_table_config(tmpl_table)

    # 收集各公司 diffs
    company_diffs = {}
    for company_dir in COMPANY_DIRS:
        data_dir = os.path.join(SCRIPT_DIR, company_dir)
        _, _, diffs = compare_order_prices(
            os.path.join(data_dir, "orders.csv"),
            os.path.join(data_dir, "orders2.csv"),
        )
        company_diffs[company_dir] = diffs or []

    all_abs = [abs(d) for diffs in company_diffs.values() for d in diffs]
    if not all_abs:
        print("  No price differences found, skipping summary.")
        return

    # 动态计算区间：1-100, 101-200, 201-300, ...
    max_val = int(max(all_abs))
    max_ceiling = ((max_val - 1) // 100 + 1) * 100
    ranges = [(i, min(i + 99, max_ceiling)) for i in range(1, max_ceiling + 1, 100)]

    # 表头和数据
    headers = ["Company"] + [f"{lo}-{hi}" for lo, hi in ranges]
    table_data = []
    for company_dir in COMPANY_DIRS:
        diffs = company_diffs[company_dir]
        row = [company_dir]
        for lo, hi in ranges:
            row.append(str(sum(1 for d in diffs if lo <= abs(d) <= hi)))
        table_data.append(row)

    # ── 封面 ─────────────────────────────────────
    slide = _clone_slide(prs, tmpl_cover)
    sh = _find_shape(slide, "Title")
    if sh:
        _set_text(sh, "Price Difference Summary")
    sh = _find_shape(slide, "Subtitle")
    if sh:
        _set_text(sh, "Statistical overview across all companies")

    # ── 分隔页 ───────────────────────────────────
    build_section(prs, tmpl_section, "Price Difference Distribution", len(table_data), len(headers))

    # ── 表格页 ───────────────────────────────────
    tbl_shape = build_table(prs, tmpl_table, headers, table_data, 1, 1,
                            "Price Difference Distribution", row_height, max_table_height,
                            header_style, data_style, header_fill, alt_fill)

    # 为公司名添加超链接，点击跳转到对应 report.pptx 的价格对比页
    tbl = tbl_shape.table
    for ri, company_dir in enumerate(COMPANY_DIRS):
        cell = tbl.cell(ri + 1, 0)
        if not cell.text_frame.paragraphs[0].runs:
            continue
        run = cell.text_frame.paragraphs[0].runs[0]

        report_path = os.path.join(SCRIPT_DIR, company_dir, "report.pptx")
        slide_num, sld_id, slide_title = _get_price_comparison_slide_info(report_path)

        if slide_num is None:
            print(f"  Warning: No price comparison slide found in {report_path}")
            continue

        # PowerPoint OOXML 格式：
        #   Relationship Target = "company/report.pptx#<sldId>,<slideNum>,<title>"
        #   hlinkClick action   = "ppaction://hlinkpres?slideindex=<slideNum>&slidetitle=<title>"
        fragment = f"{sld_id},{slide_num},{slide_title}"
        run.hyperlink.address = f"{company_dir}/report.pptx#{fragment}"

        hlinkClick = run._r.find('.//' + qn('a:hlinkClick'))
        if hlinkClick is not None:
            action = (f"ppaction://hlinkpres?slideindex={slide_num}"
                      f"&slidetitle={slide_title}")
            hlinkClick.set('action', action)

        print(f"  {company_dir}: hyperlink -> slide {slide_num} (sldId={sld_id})")

    # 删除模板页
    for idx in [2, 1, 0]:
        rId = prs.slides._sldIdLst[idx].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

    prs.save(output_path)
    print(f"\nDone -> {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def main():
    for company_dir in COMPANY_DIRS:
        print(f"\n=== {company_dir} ===")
        generate_company_report(company_dir)

    print("\n=== Summary ===")
    generate_summary_report()
    print("\nAll reports generated.")


if __name__ == "__main__":
    main()
