"""
create_template.py
生成 report_template.pptx 模板文件。

模板包含 3 张预设幻灯片，各自带有装饰元素（色条等）和命名的占位文本框：
  Slide 0 — COVER       : 左侧蓝色竖条 + 标题/副标题/分隔线/内容列表区
  Slide 1 — SECTION     : 顶部蓝色横条 + 居中标题/描述区
  Slide 2 — DATA_TABLE  : 顶部细蓝色条 + 小标题/页码区

数据脚本加载后，删除这 3 张模板页，但基于它们的布局信息逐张复制生成新幻灯片。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from pptx.oxml.ns import qn

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

PRISM_BLUE = RGBColor(0x00, 0x70, 0xC0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "report_template.pptx")


def _bar(slide, left, top, width, height, color=PRISM_BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _textbox(slide, name, left, top, width, height, text="", size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tx.name = name
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tx


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank = prs.slide_layouts[6]  # blank layout

    # ── Slide 0: COVER ────────────────────────────────
    s0 = prs.slides.add_slide(blank)
    _bar(s0, 0, 0, Inches(0.35), SLIDE_HEIGHT)                              # left accent
    _bar(s0, Inches(1.2), Inches(3.9), Inches(3), Pt(3))                    # divider
    _textbox(s0, "Title",       Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
             "REPORT TITLE", size=44, bold=True, color=PRISM_BLUE)
    _textbox(s0, "Subtitle",    Inches(1.2), Inches(3.1), Inches(10), Inches(0.6),
             "Subtitle goes here", size=22, color=MEDIUM_GRAY)
    _textbox(s0, "ContentList", Inches(1.2), Inches(4.3), Inches(10), Inches(2.5),
             "", size=14, color=MEDIUM_GRAY)

    # ── Slide 1: SECTION ──────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _bar(s1, 0, 0, SLIDE_WIDTH, Inches(0.15))
    _textbox(s1, "SectionTitle",  Inches(1), Inches(2.4), Inches(11), Inches(1.0),
             "SECTION TITLE", size=36, bold=True, color=PRISM_BLUE, align=PP_ALIGN.CENTER)
    _textbox(s1, "DimensionInfo", Inches(1), Inches(3.6), Inches(11), Inches(0.6),
             "NxN", size=20, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 2: DATA_TABLE ───────────────────────────
    # 包含一个示例表格（1行表头 + 10行数据），用户可增减行数来控制每页数据量
    s2 = prs.slides.add_slide(blank)
    _bar(s2, 0, 0, SLIDE_WIDTH, Inches(0.08))
    _textbox(s2, "SectionTitle", Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
             "Section", size=14, bold=True, color=PRISM_BLUE)
    _textbox(s2, "PageInfo",    Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.5),
             "Page X / Y", size=12, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT)

    # 示例表格: 行数决定了最终报告每页展示的数据行数
    # 默认 11 行 = 1 行表头 + 10 行数据
    SAMPLE_ROWS = 11  # 1 header + 10 data
    SAMPLE_COLS = 5
    tbl_shape = s2.shapes.add_table(
        SAMPLE_ROWS, SAMPLE_COLS,
        Inches(0.5), Inches(0.85),
        Inches(12.3), Inches(6.3),
    )
    tbl_shape.name = "SampleTable"

    # 给示例表格填入提示文字
    tbl = tbl_shape.table
    header_texts = ["Column 1", "Column 2", "Column 3", "Column 4", "Column 5"]
    for ci, h in enumerate(header_texts):
        cell = tbl.cell(0, ci)
        cell.text = h
    for ri in range(1, SAMPLE_ROWS):
        for ci in range(SAMPLE_COLS):
            tbl.cell(ri, ci).text = f"data"

    prs.save(TEMPLATE_PATH)
    print(f"Template saved -> {TEMPLATE_PATH}")
    print(f"Template slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
